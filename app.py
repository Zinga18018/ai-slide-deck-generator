import os
import re
import logging
import os
from datetime import datetime
from dotenv import load_dotenv
from flask import Flask, render_template, request, send_file, jsonify, flash, redirect, url_for
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from werkzeug.utils import secure_filename
import google.generativeai as genai
import json
import tempfile

# Load environment variables from .env file
load_dotenv()

app = Flask(__name__)
app.secret_key = 'neuralslide-ai-secret-key-2024'  # Enhanced secret key

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Configure Gemini API
API_KEY = os.getenv('GEMINI_API_KEY', 'your-api-key-here')  # Get from environment variable
if API_KEY == 'your-api-key-here' or not API_KEY:
    logger.warning("⚠️  GEMINI_API_KEY not set! Please set your API key as an environment variable.")
    logger.info("💡 Get your API key from: https://makersuite.google.com/app/apikey")
    logger.info("💡 Set it with: set GEMINI_API_KEY=your_actual_api_key (Windows)")
    logger.info("💡 Or create a .env file with: GEMINI_API_KEY=your_actual_api_key")

genai.configure(api_key=API_KEY)
model = genai.GenerativeModel('gemini-1.5-flash')

# Enhanced configuration
CONFIG = {
    'MAX_SLIDES': 25,
    'MIN_SLIDES': 3,
    'UPLOAD_FOLDER': 'static',
    'ALLOWED_EXTENSIONS': {'txt', 'pdf', 'docx'},
    'PRESENTATION_STYLES': {
        'professional': {
            'primary_color': RGBColor(31, 73, 125),
            'secondary_color': RGBColor(68, 114, 196),
            'accent_color': RGBColor(255, 255, 255),
            'text_color': RGBColor(64, 64, 64)
        },
        'creative': {
            'primary_color': RGBColor(255, 107, 107),
            'secondary_color': RGBColor(78, 205, 196),
            'accent_color': RGBColor(255, 255, 255),
            'text_color': RGBColor(64, 64, 64)
        },
        'academic': {
            'primary_color': RGBColor(68, 68, 68),
            'secondary_color': RGBColor(128, 128, 128),
            'accent_color': RGBColor(255, 255, 255),
            'text_color': RGBColor(32, 32, 32)
        },
        'business': {
            'primary_color': RGBColor(0, 32, 96),
            'secondary_color': RGBColor(0, 112, 192),
            'accent_color': RGBColor(255, 255, 255),
            'text_color': RGBColor(64, 64, 64)
        },
        'minimal': {
            'primary_color': RGBColor(64, 64, 64),
            'secondary_color': RGBColor(128, 128, 128),
            'accent_color': RGBColor(255, 255, 255),
            'text_color': RGBColor(96, 96, 96)
        }
    }
}

# Ensure static directory exists
os.makedirs(CONFIG['UPLOAD_FOLDER'], exist_ok=True)

@app.route('/')
def index():
    """Render the main page with the text input form."""
    return render_template('index.html')

def generate_ai_content(user_text, style='professional', slide_count='8-12', include_charts=False, speaker_notes=False):
    """
    Generate advanced structured presentation content using Gemini AI
    """
    # Parse slide count
    min_slides, max_slides = map(int, slide_count.split('-'))
    
    # Enhanced prompt with style and options
    prompt = f"""
    You are an expert presentation designer and content strategist. Create a comprehensive, {style} slide deck from the following content.
    
    Content: {user_text}
    
    Requirements:
    - Style: {style.title()} (adjust tone and language accordingly)
    - Number of slides: {min_slides}-{max_slides} slides
    - Include charts/graphs: {include_charts}
    - Speaker notes: {speaker_notes}
    
    Please generate a JSON response with this exact structure:
    {{
        "title": "Compelling presentation title",
        "subtitle": "Brief subtitle or tagline",
        "author": "Generated by NeuralSlide AI",
        "slides": [
            {{
                "title": "Slide title",
                "content": ["Key point 1", "Key point 2", "Key point 3"],
                "slide_type": "content|intro|conclusion|data|comparison",
                "chart_data": {{
                    "type": "bar|line|pie|column",
                    "title": "Chart title",
                    "categories": ["Cat1", "Cat2", "Cat3"],
                    "values": [10, 20, 30]
                }} // Only if include_charts is true and relevant
                {f', "speaker_notes": "Detailed speaker notes for this slide"' if speaker_notes else ''}
            }}
        ]
    }}
    
    Content Guidelines:
    - Create engaging, actionable content
    - Use data-driven insights when possible
    - Include real-world examples and case studies
    - Ensure logical narrative flow
    - Make each slide self-contained but connected
    - Use compelling headlines and clear bullet points
    - For {style} style: {'Use formal language, data focus, and professional terminology' if style == 'professional' else 'Use creative language, engaging visuals, and innovative approaches' if style == 'creative' else 'Use academic language, research focus, and scholarly approach' if style == 'academic' else 'Use business language, ROI focus, and strategic thinking' if style == 'business' else 'Use clean language, essential points only, and minimal text'}
    
    Slide Types:
    - intro: Introduction/title slides
    - content: Main content slides
    - data: Data-heavy slides with charts
    - comparison: Comparison or vs. slides
    - conclusion: Summary/conclusion slides
    
    Return ONLY the JSON, no additional text or formatting.
    """
    
    try:
        logger.info(f"Generating AI content with style: {style}, slides: {slide_count}")
        response = model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        logger.error(f"Error generating AI content: {e}")
        return None

@app.route('/generate', methods=['POST'])
def generate_presentation():
    try:
        # Get user input and settings
        user_text = request.form.get('text_input', '').strip()
        style = request.form.get('presentation_style', 'professional')
        slide_count = request.form.get('slide_count', '8-12')
        include_charts = request.form.get('include_charts') == 'on'
        speaker_notes = request.form.get('speaker_notes') == 'on'
        
        logger.info(f"Processing request: style={style}, slides={slide_count}, charts={include_charts}, notes={speaker_notes}")
        
        if not user_text:
            flash('Please enter some text to generate a presentation.', 'error')
            return redirect(url_for('index'))
        
        if len(user_text) < 50:
            flash('Please provide more detailed content (at least 50 characters) for better results.', 'error')
            return redirect(url_for('index'))
        
        # Generate AI content with advanced parameters
        ai_response = generate_ai_content(user_text, style, slide_count, include_charts, speaker_notes)
        
        if not ai_response:
            flash('Failed to generate AI content. Please check your API key and try again.', 'error')
            return redirect(url_for('index'))
        
        # Parse JSON response with enhanced error handling
        try:
            # Clean the response to extract JSON
            ai_response = ai_response.strip()
            logger.info(f"Raw AI response length: {len(ai_response)}")
            
            # Remove markdown code blocks if present
            if ai_response.startswith('```'):
                lines = ai_response.split('\n')
                # Handle different markdown formats
                if len(lines) > 2:
                    ai_response = '\n'.join(lines[1:-1])
                else:
                    ai_response = ai_response.replace('```', '')
            
            # Remove any leading/trailing non-JSON text
            ai_response = ai_response.strip()
            
            # Find JSON boundaries more robustly
            json_start = ai_response.find('{')
            json_end = ai_response.rfind('}') + 1
            
            if json_start == -1 or json_end == 0:
                logger.error(f"No JSON boundaries found in response: {ai_response[:200]}...")
                raise ValueError("No JSON found in response")
            
            json_str = ai_response[json_start:json_end]
            logger.info(f"Extracted JSON length: {len(json_str)}")
            
            # Enhanced JSON cleaning
            # Remove trailing commas
            json_str = re.sub(r',\s*}', '}', json_str)
            json_str = re.sub(r',\s*]', ']', json_str)
            # Remove comments
            json_str = re.sub(r'//.*?\n', '\n', json_str)
            json_str = re.sub(r'/\*.*?\*/', '', json_str, flags=re.DOTALL)
            # Fix common JSON issues
            json_str = re.sub(r'\n\s*\n', '\n', json_str)  # Remove extra newlines
            json_str = json_str.replace('\"', '"')  # Fix escaped quotes
            
            # Try to parse JSON
            presentation_data = json.loads(json_str)
            logger.info("JSON parsed successfully")
            
        except (json.JSONDecodeError, ValueError) as e:
            logger.error(f"JSON parsing error: {e}")
            logger.error(f"AI Response (first 1000 chars): {ai_response[:1000]}")
            if 'json_str' in locals():
                logger.error(f"Cleaned JSON (first 500 chars): {json_str[:500]}")
            
            # Try alternative parsing methods
            try:
                # Attempt to fix common JSON issues and retry
                if 'json_str' in locals():
                    # Try to fix incomplete JSON
                    if json_str.count('{') > json_str.count('}'):
                        json_str += '}' * (json_str.count('{') - json_str.count('}'))
                    if json_str.count('[') > json_str.count(']'):
                        json_str += ']' * (json_str.count('[') - json_str.count(']'))
                    
                    presentation_data = json.loads(json_str)
                    logger.info("JSON parsed successfully after repair")
                else:
                    raise e
            except:
                flash('Failed to parse AI response. The AI response format was invalid. Please try again with different content or check if the API is working properly.', 'error')
                return redirect(url_for('index'))
        
        # Validate the structure
        if not all(key in presentation_data for key in ['title', 'slides']):
            flash('Invalid AI response format. Please try again with different content.', 'error')
            return redirect(url_for('index'))
        
        if not presentation_data['slides'] or len(presentation_data['slides']) == 0:
            flash('No slides were generated. Please provide more detailed content.', 'error')
            return redirect(url_for('index'))
        
        # Create enhanced PowerPoint presentation
        filename = create_enhanced_powerpoint(presentation_data, style, include_charts, speaker_notes)
        
        if filename:
            flash(f'🎉 Presentation "{presentation_data["title"]}" generated successfully with {len(presentation_data["slides"])} slides!', 'success')
            return render_template('download.html', 
                                 filename=filename, 
                                 title=presentation_data['title'],
                                 slide_count=len(presentation_data['slides']),
                                 style=style)
        else:
            flash('Failed to create presentation file. Please try again.', 'error')
            return redirect(url_for('index'))
            
    except Exception as e:
        logger.error(f"Error in generate_presentation: {e}")
        flash('An unexpected error occurred. Please try again or contact support.', 'error')
        return redirect(url_for('index'))

@app.route('/download/<filename>')
def download_file(filename):
    """Download the generated PowerPoint file."""
    try:
        filepath = os.path.join('static', filename)
        if os.path.exists(filepath):
            return send_file(filepath, as_attachment=True, download_name=filename)
        else:
            flash('File not found.', 'error')
            return redirect(url_for('index'))
    except Exception as e:
        flash(f'Error downloading file: {str(e)}', 'error')
        return redirect(url_for('index'))

def create_enhanced_powerpoint(presentation_data, style='professional', include_charts=False, speaker_notes=False):
    """Create an enhanced PowerPoint presentation with styling and advanced features."""
    try:
        prs = Presentation()
        style_config = CONFIG['PRESENTATION_STYLES'].get(style, CONFIG['PRESENTATION_STYLES']['professional'])
        
        # Add title slide with custom styling
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Style title slide
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = presentation_data.get('title', 'AI Generated Presentation')
        subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')} | Style: {style.title()}"
        
        # Apply title styling
        if title.text_frame:
            for paragraph in title.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(*style_config['primary_color'])
                    run.font.size = Pt(44)
                    run.font.bold = True
        
        if subtitle.text_frame:
            for paragraph in subtitle.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(*style_config['secondary_color'])
                    run.font.size = Pt(18)
        
        # Add content slides with enhanced features
        chart_data_available = []
        
        for i, slide_data in enumerate(presentation_data['slides']):
            # Choose layout based on content type
            if include_charts and slide_data.get('chart_data'):
                slide_layout = prs.slide_layouts[5]  # Blank layout for custom chart
            else:
                slide_layout = prs.slide_layouts[1]  # Title and Content
            
            slide = prs.slides.add_slide(slide_layout)
            
            # Set slide title with styling
            title_shape = slide.shapes.title
            title_shape.text = slide_data['title']
            
            if title_shape.text_frame:
                for paragraph in title_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(*style_config['primary_color'])
                        run.font.size = Pt(32)
                        run.font.bold = True
            
            # Add content based on type
            if include_charts and slide_data.get('chart_data'):
                # Create chart slide
                try:
                    chart_data = CategoryChartData()
                    chart_info = slide_data['chart_data']
                    
                    chart_data.categories = chart_info.get('categories', ['Category 1', 'Category 2', 'Category 3'])
                    chart_data.add_series('Data', chart_info.get('values', [10, 20, 30]))
                    
                    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(5)
                    chart = slide.shapes.add_chart(
                        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
                    ).chart
                    
                    chart.has_legend = True
                    chart_data_available.append(i + 1)
                    
                except Exception as chart_error:
                    logger.warning(f"Failed to create chart for slide {i+1}: {chart_error}")
                    # Fall back to regular content
                    _add_regular_content(slide, slide_data, style_config)
            else:
                _add_regular_content(slide, slide_data, style_config)
            
            # Add speaker notes if requested
            if speaker_notes and slide_data.get('speaker_notes'):
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                notes_text_frame.text = slide_data['speaker_notes']
        
        # Add summary slide if there are many slides
        if len(presentation_data['slides']) > 5:
            summary_layout = prs.slide_layouts[1]
            summary_slide = prs.slides.add_slide(summary_layout)
            
            summary_title = summary_slide.shapes.title
            summary_title.text = "Summary & Key Takeaways"
            
            if summary_title.text_frame:
                for paragraph in summary_title.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(*style_config['primary_color'])
                        run.font.size = Pt(32)
                        run.font.bold = True
            
            # Add summary points
            content_shape = summary_slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            summary_points = presentation_data.get('summary', [
                "Key insights from the presentation",
                "Main conclusions and recommendations",
                "Next steps and action items"
            ])
            
            for i, point in enumerate(summary_points[:5]):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = str(point)
                p.level = 0
                
                for run in p.runs:
                    run.font.color.rgb = RGBColor(*style_config['text_color'])
                    run.font.size = Pt(18)
        
        # Save the presentation with descriptive filename
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        safe_title = re.sub(r'[^\w\s-]', '', presentation_data.get('title', 'presentation'))[:30]
        safe_title = re.sub(r'[-\s]+', '_', safe_title)
        filename = f'{safe_title}_{style}_{timestamp}.pptx'
        filepath = os.path.join('static', filename)
        
        prs.save(filepath)
        
        logger.info(f"Created presentation: {filename} with {len(presentation_data['slides'])} slides")
        if chart_data_available:
            logger.info(f"Charts added to slides: {chart_data_available}")
        
        return filename
        
    except Exception as e:
        logger.error(f"Error creating enhanced PowerPoint: {e}")
        return None

def _add_regular_content(slide, slide_data, style_config):
    """Helper function to add regular text content to a slide."""
    try:
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # Add bullet points or content
        content = slide_data.get('content', [])
        if isinstance(content, list):
            for i, bullet_point in enumerate(content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = str(bullet_point)
                p.level = 0
                
                for run in p.runs:
                    run.font.color.rgb = RGBColor(*style_config['text_color'])
                    run.font.size = Pt(18)
        else:
            text_frame.paragraphs[0].text = str(content)
            for run in text_frame.paragraphs[0].runs:
                run.font.color.rgb = RGBColor(*style_config['text_color'])
                run.font.size = Pt(18)
                
    except Exception as e:
        logger.warning(f"Error adding content to slide: {e}")

@app.errorhandler(500)
def internal_error(error):
    return render_template('error.html', error="Internal server error occurred. Please try again."), 500

@app.errorhandler(404)
def not_found(error):
    return render_template('error.html', error="Page not found."), 404

@app.route('/health')
def health_check():
    return {'status': 'healthy', 'timestamp': datetime.now().isoformat()}

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    debug_mode = os.environ.get('FLASK_ENV') != 'production'
    app.run(debug=debug_mode, host='0.0.0.0', port=port)